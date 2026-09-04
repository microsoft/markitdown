"""Tests for PPTX chart conversion.

Scatter and bubble charts have no category axis, so their data must be read
point-by-point. Previously ``_convert_chart_to_markdown`` iterated over
``chart.plots[0].categories`` (empty for these chart types) and silently
dropped the entire series. These tests assert the data points are preserved.
"""
import io

from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE

from markitdown import MarkItDown


def _build_xy_pptx(bubble):
    """Build an in-memory PPTX containing a single scatter or bubble chart."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if bubble:
        chart_data = BubbleChartData()
        series = chart_data.add_series("Series 1")
        series.add_data_point(1.0, 2.0, 10)
        series.add_data_point(2.0, 3.0, 20)
        series.add_data_point(3.0, 4.0, 30)
        chart_type = XL_CHART_TYPE.BUBBLE
    else:
        chart_data = XyChartData()
        series = chart_data.add_series("Series 1")
        series.add_data_point(1.0, 2.0)
        series.add_data_point(2.0, 3.0)
        series.add_data_point(3.0, 4.0)
        series2 = chart_data.add_series("Series 2")
        series2.add_data_point(1.0, 5.0)
        series2.add_data_point(2.0, 6.0)
        series2.add_data_point(3.0, 7.0)
        chart_type = XL_CHART_TYPE.XY_SCATTER
    slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1), Inches(5), Inches(4), chart_data
    )
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def test_scatter_chart_preserves_data_points() -> None:
    # Regression: scatter series used to be dropped entirely (only the header
    # was emitted) because the chart has no category axis.
    md = (
        MarkItDown()
        .convert_stream(_build_xy_pptx(bubble=False), file_extension=".pptx")
        .markdown
    )
    assert "| X | Series 1 | Series 2 |" in md
    assert "| 1.0 | 2.0 | 5.0 |" in md
    assert "| 3.0 | 4.0 | 7.0 |" in md


def test_bubble_chart_preserves_data_points_and_sizes() -> None:
    # Bubble charts additionally carry a per-point size, which must appear in a
    # dedicated "(size)" column.
    md = (
        MarkItDown()
        .convert_stream(_build_xy_pptx(bubble=True), file_extension=".pptx")
        .markdown
    )
    assert "| X | Series 1 | Series 1 (size) |" in md
    assert "| 1.0 | 2.0 | 10.0 |" in md
    assert "| 3.0 | 4.0 | 30.0 |" in md


if __name__ == "__main__":
    test_scatter_chart_preserves_data_points()
    test_bubble_chart_preserves_data_points_and_sizes()
    print("All chart tests passed!")
