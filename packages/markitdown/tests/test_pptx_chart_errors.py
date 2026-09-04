#!/usr/bin/env python3 -m pytest
"""Tests for chart reads that fail inside PptxConverter.

``_convert_chart_to_markdown`` wraps its body in two handlers that both return
the same ``[unsupported chart]`` placeholder, but the ``ValueError`` handler
returned it only when the message contained "unsupported plot type". Any other
``ValueError`` reached the end of the handler and fell off the function, so the
method returned ``None``.

The caller does ``md_content += self._convert_chart_to_markdown(shape.chart)``
unguarded, so that ``None`` raises ``TypeError: can only concatenate str (not
"NoneType") to str``. ``MarkItDown._convert`` catches it and, with no other
converter accepting a .pptx, re-raises as ``FileConversionException`` -- one
unreadable chart therefore loses the entire presentation rather than a single
shape.
"""

import io
from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from markitdown import MarkItDown
from markitdown._stream_info import StreamInfo
from markitdown.converters._pptx_converter import PptxConverter

PLACEHOLDER = "[unsupported chart]"

# A ValueError python-pptx can raise that is not the special-cased one.
OTHER_VALUE_ERROR = "chart has no categories"


def _deck_with_chart() -> bytes:
    """A minimal one-slide deck holding a plain clustered column chart."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("S1", (1.0, 2.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _convert(deck: bytes) -> str:
    return (
        MarkItDown()
        .convert_stream(io.BytesIO(deck), stream_info=StreamInfo(extension=".pptx"))
        .markdown
    )


class _ChartRaising:
    """Stands in for a chart whose first access raises."""

    def __init__(self, message: str):
        self._message = message

    @property
    def has_title(self):
        raise ValueError(self._message)


@pytest.mark.parametrize("message", [OTHER_VALUE_ERROR, "unsupported plot type xlFoo"])
def test_chart_read_failure_returns_a_string(message: str) -> None:
    """The method must never hand the caller a None to concatenate."""
    result = PptxConverter()._convert_chart_to_markdown(_ChartRaising(message))

    assert isinstance(result, str)
    assert PLACEHOLDER in result


def test_deck_survives_an_unreadable_chart() -> None:
    """A chart that cannot be read costs the chart, not the presentation."""
    deck = _deck_with_chart()

    def _raise(self):
        raise ValueError(OTHER_VALUE_ERROR)

    with patch.object(Chart, "plots", property(_raise)):
        markdown = _convert(deck)

    assert PLACEHOLDER in markdown


def test_readable_chart_is_still_converted() -> None:
    """Control: an ordinary chart must keep converting to a table."""
    markdown = _convert(_deck_with_chart())

    assert "### Chart" in markdown
    assert PLACEHOLDER not in markdown
    assert "| Category | S1 |" in markdown


if __name__ == "__main__":
    test_chart_read_failure_returns_a_string(OTHER_VALUE_ERROR)
    test_chart_read_failure_returns_a_string("unsupported plot type xlFoo")
    test_deck_survives_an_unreadable_chart()
    test_readable_chart_is_still_converted()
    print("All tests passed!")
