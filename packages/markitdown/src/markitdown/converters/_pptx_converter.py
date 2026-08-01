import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any
from operator import attrgetter

from ._html_converter import HtmlConverter
from ._llm_caption import llm_caption
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._stream_info import StreamInfo
from .._exceptions import MissingDependencyException, MISSING_DEPENDENCY_MESSAGE

# Try loading optional (but in this case, required) dependencies
# Save reporting of any exceptions for later
_dependency_exc_info = None
try:
    import pptx
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    # Preserve the error and stack trace for later
    _dependency_exc_info = sys.exc_info()


ACCEPTED_MIME_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]

ACCEPTED_FILE_EXTENSIONS = [".pptx"]

# Chart types drawn on an X/Y (value) axis instead of a category axis. They
# have no ``chart.plots[0].categories``, so their series data must be read
# point-by-point rather than row-by-row over categories.
XY_CHART_TYPES = ()
BUBBLE_CHART_TYPES = ()
if _dependency_exc_info is None:
    XY_CHART_TYPES = (
        XL_CHART_TYPE.XY_SCATTER,
        XL_CHART_TYPE.XY_SCATTER_LINES,
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
        XL_CHART_TYPE.BUBBLE,
        XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT,
    )
    BUBBLE_CHART_TYPES = (XL_CHART_TYPE.BUBBLE, XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT)


def _chart_to_markdown_table(data):
    """Render a list of rows (first row is the header) as a Markdown table."""
    if not data:
        return ""
    rows = ["| " + " | ".join(map(str, row)) + " |" for row in data]
    header = rows[0]
    separator = "|" + "|".join(["---"] * len(data[0])) + "|"
    return "\n".join([header, separator] + rows[1:])


def _xy_point_x(series, idx):
    """X value of point *idx* in an XY (scatter) series.

    python-pptx only exposes the Y values publicly (``series.values``); the X
    values live on the internal ``c:xVal`` element, so we read them directly.
    """
    x_val = series._element.xVal
    if x_val is None or idx >= x_val.ptCount_val:
        return None
    return x_val.pt_v(idx)


def _xy_point_size(series, idx):
    """Bubble size of point *idx* in a bubble series."""
    size = series._element.bubbleSize
    if size is None or idx >= size.ptCount_val:
        return None
    return size.pt_v(idx)


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self):
        super().__init__()
        self._html_converter = HtmlConverter()

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension in ACCEPTED_FILE_EXTENSIONS:
            return True

        for prefix in ACCEPTED_MIME_TYPE_PREFIXES:
            if mimetype.startswith(prefix):
                return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:
        # Check the dependencies
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[
                1
            ].with_traceback(  # type: ignore[union-attr]
                _dependency_exc_info[2]
            )

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    llm_description = ""
                    alt_text = ""

                    # Resolve the image blob, handling SVG images that lack a
                    # rasterized fallback (shape.image raises in that case).
                    (
                        image_blob,
                        image_content_type,
                        image_filename,
                    ) = self._get_image_info(shape)

                    # Potentially generate a description using an LLM
                    llm_client = kwargs.get("llm_client")
                    llm_model = kwargs.get("llm_model")
                    if (
                        llm_client is not None
                        and llm_model is not None
                        and image_blob is not None
                    ):
                        # Prepare a file_stream and stream_info for the image data
                        image_extension = None
                        if image_filename:
                            image_extension = os.path.splitext(image_filename)[1]
                        image_stream_info = StreamInfo(
                            mimetype=image_content_type,
                            extension=image_extension,
                            filename=image_filename,
                        )

                        image_stream = io.BytesIO(image_blob)

                        # Caption the image
                        try:
                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=llm_model,
                                prompt=kwargs.get("llm_prompt"),
                            )
                        except Exception:
                            # Unable to generate a description
                            pass

                    # Also grab any description embedded in the deck
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        # Unable to get alt text
                        pass

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([llm_description, alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # If keep_data_uris is True, use base64 encoding for images
                    if kwargs.get("keep_data_uris", False) and image_blob is not None:
                        content_type = image_content_type or "image/png"
                        b64_string = base64.b64encode(image_blob).decode("utf-8")
                        md_content += f"\n![{alt_text}](data:{content_type};base64,{b64_string})\n"
                    else:
                        # A placeholder name
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        md_content += "\n![" + alt_text + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _find_svg_blip_part(self, shape):
        """Return the image part referenced by an ``<asvg:svgBlip>``, if any.

        PowerPoint stores SVG pictures as a blip whose main ``r:embed`` points
        to a rasterized PNG fallback, plus an ``<asvg:svgBlip>`` extension
        pointing to the SVG. When there is no raster fallback the ``<a:blip>``
        has no ``r:embed`` at all, so python-pptx's ``shape.image`` fails. This
        resolves the SVG part directly from the ``svgBlip`` extension.
        """
        try:
            nsmap = {
                "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
            }
            r_embed = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            for svg_blip in shape._element.findall(".//asvg:svgBlip", nsmap):
                embed_rid = svg_blip.get(r_embed)
                if not embed_rid:
                    continue
                return shape.part.related_part(embed_rid)
        except Exception:
            pass
        return None

    def _get_image_info(self, shape):
        """Return (blob, content_type, filename) for a picture shape.

        Handles SVG images that lack a rasterized fallback. In that case
        ``shape.image`` raises ``ValueError("no embedded image")`` because the
        ``<a:blip>`` has no ``r:embed`` attribute (only an ``<asvg:svgBlip>``
        extension). We fall back to resolving the SVG blip directly.
        """
        try:
            image = shape.image
            return image.blob, image.content_type, image.filename
        except Exception:
            pass

        # Fall back to an embedded SVG blip (image without a raster fallback)
        part = self._find_svg_blip_part(shape)
        if part is not None:
            try:
                filename = os.path.basename(getattr(part, "partname", "") or "") or None
                return part.blob, "image/svg+xml", filename
            except Exception:
                pass

        return None, None, None

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            # ``shape.image`` can raise (e.g. ValueError "no embedded image")
            # for SVG placeholders without a raster fallback, so guard against
            # any exception rather than relying on hasattr (which only swallows
            # AttributeError).
            try:
                if shape.image is not None:
                    return True
            except Exception:
                # Still a picture if it carries an embedded SVG blip.
                if self._find_svg_blip_part(shape) is not None:
                    return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"

            series_list = list(chart.series)
            series_names = [s.name for s in series_list]
            # Materialize each series' Y values once. Accessing series.values[idx]
            # inside the nested loop is O(n^2) in python-pptx (each lookup does an
            # XPath scan of all points), which is extremely slow on large charts.
            series_values = [list(s.values) for s in series_list]

            # Scatter and bubble charts have no category axis: each data point
            # carries its own X coordinate, so chart.plots[0].categories is empty.
            # Iterating over categories (as category charts do) would emit zero
            # data rows and silently drop the entire series.
            if chart.chart_type in XY_CHART_TYPES:
                is_bubble = chart.chart_type in BUBBLE_CHART_TYPES
                header = ["X"]
                for name in series_names:
                    header.append(name)
                    if is_bubble:
                        header.append(f"{name} (size)")
                data = [header]
                num_points = max((len(sv) for sv in series_values), default=0)
                for idx in range(num_points):
                    # X is taken from the first series; scatter charts commonly
                    # share the X axis across series.
                    row = [_xy_point_x(series_list[0], idx)]
                    for series, sv in zip(series_list, series_values):
                        row.append(sv[idx] if idx < len(sv) else None)
                        if is_bubble:
                            row.append(_xy_point_size(series, idx))
                    data.append(row)
                return md + _chart_to_markdown_table(data)

            category_names = [c.label for c in chart.plots[0].categories]
            data = [["Category"] + series_names]
            for idx, category in enumerate(category_names):
                row = [category]
                for sv in series_values:
                    row.append(sv[idx] if idx < len(sv) else None)
                data.append(row)
            return md + _chart_to_markdown_table(data)
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
