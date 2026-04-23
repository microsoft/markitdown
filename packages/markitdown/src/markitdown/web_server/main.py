import os
import re
import json
import shutil
import uuid
import asyncio
import warnings
import zipfile
import io
import base64
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from pathlib import Path
from dataclasses import dataclass, field, asdict
from urllib.parse import unquote, urlparse

warnings.filterwarnings("ignore", message="Couldn't find ffmpeg or avconv")

from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks, Form
from fastapi.responses import HTMLResponse, JSONResponse, FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from markitdown import MarkItDown
from markitdown._exceptions import (
    MarkItDownException,
    FileConversionException,
    UnsupportedFormatException,
)


app = FastAPI(
    title="MarkItDown Web UI",
    description="A lightweight web interface for MarkItDown file conversion",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


UPLOAD_DIR = Path("uploads")
RESULT_DIR = Path("results")
IMAGES_DIR = Path("results/images")
HISTORY_FILE = Path("results/conversion_history.json")
UPLOAD_DIR.mkdir(exist_ok=True)
RESULT_DIR.mkdir(exist_ok=True)
IMAGES_DIR.mkdir(exist_ok=True)


conversion_tasks: Dict[str, Any] = {}


@dataclass
class ConversionOptions:
    output_format: str = "standard"
    extract_images: bool = False
    optimize_tables: bool = False
    optimize_code_blocks: bool = False


@dataclass
class ConversionProgress:
    task_id: str
    filename: str
    status: str
    progress: int
    result: Optional[str] = None
    error_message: Optional[str] = None
    created_at: datetime = field(default_factory=datetime.now)
    completed_at: Optional[datetime] = None
    options: ConversionOptions = field(default_factory=ConversionOptions)
    images: List[str] = field(default_factory=list)
    original_file_size: int = 0


class ConversionResponse(BaseModel):
    task_id: str
    filename: str
    status: str
    progress: int


class ProgressResponse(BaseModel):
    task_id: str
    filename: str
    status: str
    progress: int
    error_message: Optional[str] = None
    images: Optional[List[str]] = None


class BatchConversionResponse(BaseModel):
    tasks: List[ConversionResponse]


class HistoryItem(BaseModel):
    task_id: str
    filename: str
    status: str
    created_at: str
    completed_at: Optional[str] = None
    output_format: str
    images_extracted: bool
    file_size: int


class ClearCacheResponse(BaseModel):
    success: bool
    message: str
    uploads_cleared: int
    results_cleared: int
    images_cleared: int


markitdown_instance = MarkItDown()


SUPPORTED_EXTENSIONS = [
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".html", ".htm", ".csv", ".txt", ".md", ".json",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp",
    ".mp3", ".wav", ".m4a", ".flac",
    ".epub", ".zip", ".msg", ".ipynb", ".xml",
]


def is_supported_format(filename: str) -> bool:
    ext = Path(filename).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


def get_supported_formats_message() -> str:
    formats = ", ".join(sorted(set(SUPPORTED_EXTENSIONS)))
    return f"Supported formats: {formats}"


def optimize_tables_format(markdown: str) -> str:
    lines = markdown.split('\n')
    optimized_lines = []
    in_table = False
    table_buffer = []
    
    for line in lines:
        if re.match(r'^\|.+\|$', line) and not re.match(r'^\|[-:| ]+\|$', line):
            if not in_table:
                in_table = True
                table_buffer = []
            table_buffer.append(line)
        elif re.match(r'^\|[-:| ]+\|$', line):
            if in_table:
                table_buffer.append(line)
        elif in_table:
            if table_buffer:
                optimized_table = format_table_properly(table_buffer)
                optimized_lines.extend(optimized_table)
            optimized_lines.append(line)
            in_table = False
            table_buffer = []
        else:
            optimized_lines.append(line)
    
    if in_table and table_buffer:
        optimized_table = format_table_properly(table_buffer)
        optimized_lines.extend(optimized_table)
    
    return '\n'.join(optimized_lines)


def format_table_properly(table_lines: List[str]) -> List[str]:
    if len(table_lines) < 2:
        return table_lines
    
    rows = []
    separator_idx = None
    
    for i, line in enumerate(table_lines):
        if re.match(r'^\|[-:| ]+\|$', line):
            separator_idx = i
            rows.append([cell.strip() for cell in line.strip('|').split('|')])
        else:
            rows.append([cell.strip() for cell in line.strip('|').split('|')])
    
    if separator_idx is None:
        return table_lines
    
    col_count = max(len(row) for row in rows)
    col_widths = [0] * col_count
    
    for i, row in enumerate(rows):
        if i == separator_idx:
            continue
        for j, cell in enumerate(row):
            if j < col_count:
                col_widths[j] = max(col_widths[j], len(cell))
    
    formatted_rows = []
    for i, row in enumerate(rows):
        if i == separator_idx:
            separator_cells = []
            for j in range(col_count):
                if j < len(row) and ':' in row[j]:
                    if row[j].startswith(':') and row[j].endswith(':'):
                        separator_cells.append(':' + '-' * max(col_widths[j] - 2, 1) + ':')
                    elif row[j].endswith(':'):
                        separator_cells.append('-' * max(col_widths[j] - 1, 1) + ':')
                    else:
                        separator_cells.append(':' + '-' * max(col_widths[j] - 1, 1))
                else:
                    separator_cells.append('-' * col_widths[j])
            formatted_rows.append('| ' + ' | '.join(separator_cells) + ' |')
        else:
            padded_cells = []
            for j in range(col_count):
                cell = row[j] if j < len(row) else ''
                padded_cells.append(cell.ljust(col_widths[j]))
            formatted_rows.append('| ' + ' | '.join(padded_cells) + ' |')
    
    return formatted_rows


def optimize_code_blocks_format(markdown: str) -> str:
    pattern = r'```(\w*)\n(.*?)```'
    matches = list(re.finditer(pattern, markdown, re.DOTALL))
    
    if not matches:
        return markdown
    
    result = markdown
    for match in reversed(matches):
        lang = match.group(1)
        code = match.group(2)
        start, end = match.span()
        
        lines = code.rstrip('\n').split('\n')
        if len(lines) > 0:
            min_indent = min(len(line) - len(line.lstrip()) for line in lines if line.strip())
            if min_indent > 0:
                lines = [line[min_indent:] if line.strip() else line for line in lines]
        
        optimized_code = '\n'.join(lines) + '\n'
        replacement = f'```{lang}\n{optimized_code}```'
        result = result[:start] + replacement + result[end:]
    
    return result


def format_to_github_markdown(markdown: str) -> str:
    result = markdown
    result = re.sub(r'^#\s+(.+)$', r'# \1', result, flags=re.MULTILINE)
    result = re.sub(r'^##\s+(.+)$', r'## \1', result, flags=re.MULTILINE)
    result = re.sub(r'^###\s+(.+)$', r'### \1', result, flags=re.MULTILINE)
    
    result = re.sub(r'\[(.+?)\]\((.+?)\)', r'[\1](\2)', result)
    
    result = re.sub(r'(\n\s*[-*+]\s+.+){2,}', lambda m: format_github_list(m.group(0)), result)
    
    return result


def format_github_list(list_text: str) -> str:
    lines = list_text.strip().split('\n')
    formatted = []
    for line in lines:
        match = re.match(r'^(\s*)[-*+]\s+(.+)$', line)
        if match:
            indent = match.group(1)
            content = match.group(2)
            formatted.append(f'{indent}- {content}')
        else:
            formatted.append(line)
    return '\n'.join(formatted) + '\n'


def sanitize_filename(filename: str) -> str:
    filename = unquote(filename)
    filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
    filename = filename.strip('. ')
    if not filename:
        filename = 'file'
    return filename


def extract_all_images_from_pptx(file_path: Path, task_images_dir: Path, task_id: str) -> Tuple[List[Dict], List[Dict]]:
    """
    从 PPTX 提取所有图片，并返回：
    - 图片信息列表（用于替换 MD 中的路径）
    - 图片引用映射（用于匹配原始占位符）
    """
    extracted_images = []
    image_references = []
    try:
        import pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = pptx.Presentation(str(file_path))
        image_counter = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_counter += 1
                        image = shape.image
                        
                        placeholder_name = re.sub(r"\W", "", shape.name) + ".jpg"
                        
                        image_filename = image.filename or f"image_{slide_num}_{image_counter}"
                        image_ext = os.path.splitext(image_filename)[1] or '.png'
                        safe_name = sanitize_filename(os.path.splitext(image_filename)[0])
                        new_filename = f"{safe_name}_{image_counter}{image_ext}"
                        new_img_path = task_images_dir / new_filename
                        
                        with open(new_img_path, 'wb') as f:
                            f.write(image.blob)
                        
                        img_info = {
                            'original_name': image_filename,
                            'saved_name': new_filename,
                            'relative_path': f"images/{new_filename}",
                            'api_path': f"/api/images/{task_id}/{new_filename}",
                            'slide_num': slide_num,
                            'index': image_counter
                        }
                        extracted_images.append(img_info)
                        
                        image_references.append({
                            'placeholder_name': placeholder_name,
                            'shape_name': shape.name,
                            'image_info': img_info
                        })
                    except Exception:
                        continue
                
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for subshape in shape.shapes:
                            if hasattr(subshape, 'shape_type') and subshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                try:
                                    image_counter += 1
                                    image = subshape.image
                                    
                                    placeholder_name = re.sub(r"\W", "", subshape.name) + ".jpg"
                                    
                                    image_filename = image.filename or f"image_{slide_num}_{image_counter}"
                                    image_ext = os.path.splitext(image_filename)[1] or '.png'
                                    safe_name = sanitize_filename(os.path.splitext(image_filename)[0])
                                    new_filename = f"{safe_name}_{image_counter}{image_ext}"
                                    new_img_path = task_images_dir / new_filename
                                    
                                    with open(new_img_path, 'wb') as f:
                                        f.write(image.blob)
                                    
                                    img_info = {
                                        'original_name': image_filename,
                                        'saved_name': new_filename,
                                        'relative_path': f"images/{new_filename}",
                                        'api_path': f"/api/images/{task_id}/{new_filename}",
                                        'slide_num': slide_num,
                                        'index': image_counter
                                    }
                                    extracted_images.append(img_info)
                                    
                                    image_references.append({
                                        'placeholder_name': placeholder_name,
                                        'shape_name': subshape.name,
                                        'image_info': img_info
                                    })
                                except Exception:
                                    continue
                    except Exception:
                        continue
    except ImportError:
        pass
    except Exception:
        pass
    
    return extracted_images, image_references


def extract_all_images_from_docx(file_path: Path, task_images_dir: Path, task_id: str) -> Tuple[List[Dict], List[Dict]]:
    """
    从 DOCX 提取所有图片
    """
    extracted_images = []
    image_references = []
    try:
        from docx import Document
        
        doc = Document(str(file_path))
        image_counter = 0
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_counter += 1
                    image_part = rel.target_part
                    image_ext = os.path.splitext(rel.target_ref)[1] or '.png'
                    image_filename = os.path.basename(rel.target_ref)
                    safe_name = sanitize_filename(os.path.splitext(image_filename)[0])
                    new_filename = f"{safe_name}_{image_counter}{image_ext}"
                    new_img_path = task_images_dir / new_filename
                    
                    with open(new_img_path, 'wb') as f:
                        f.write(image_part.blob)
                    
                    img_info = {
                        'original_name': image_filename,
                        'saved_name': new_filename,
                        'relative_path': f"images/{new_filename}",
                        'api_path': f"/api/images/{task_id}/{new_filename}",
                        'index': image_counter
                    }
                    extracted_images.append(img_info)
                    
                    image_references.append({
                        'placeholder_name': image_filename,
                        'image_info': img_info
                    })
                except Exception:
                    continue
    except ImportError:
        pass
    except Exception:
        pass
    
    return extracted_images, image_references


def extract_base64_images(markdown: str, task_images_dir: Path, task_id: str) -> Tuple[str, List[Dict]]:
    """
    从 Markdown 中提取 base64 编码的图片
    """
    extracted_images = []
    result = markdown
    image_counter = 0
    
    data_uri_pattern = r'!\[([^\]]*)\]\((data:image/[^;]+;base64,([^)]+))\)'
    
    matches = list(re.finditer(data_uri_pattern, markdown))
    
    for match in reversed(matches):
        try:
            alt_text = match.group(1)
            full_data_uri = match.group(2)
            base64_data = match.group(3)
            start, end = match.span()
            
            mime_type = full_data_uri.split(';')[0].replace('data:', '')
            ext_map = {
                'image/png': '.png',
                'image/jpeg': '.jpg',
                'image/jpg': '.jpg',
                'image/gif': '.gif',
                'image/bmp': '.bmp',
                'image/webp': '.webp',
            }
            image_ext = ext_map.get(mime_type, '.png')
            
            image_counter += 1
            safe_name = sanitize_filename(alt_text) if alt_text else 'image'
            new_filename = f"{safe_name}_{image_counter}{image_ext}"
            new_img_path = task_images_dir / new_filename
            
            image_data = base64.b64decode(base64_data)
            with open(new_img_path, 'wb') as f:
                f.write(image_data)
            
            img_info = {
                'original_name': f"base64_image_{image_counter}",
                'saved_name': new_filename,
                'relative_path': f"images/{new_filename}",
                'api_path': f"/api/images/{task_id}/{new_filename}",
                'index': image_counter
            }
            extracted_images.append(img_info)
            
            replacement = f"![{alt_text}]({img_info['api_path']})"
            result = result[:start] + replacement + result[end:]
            
        except Exception:
            continue
    
    return result, extracted_images


def find_all_image_references(markdown: str) -> List[Dict]:
    """
    查找 Markdown 中的所有图片引用
    返回格式：[{ 'alt': '...', 'path': '...', 'start': index, 'end': index }]
    """
    img_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
    references = []
    
    for match in re.finditer(img_pattern, markdown):
        alt_text = match.group(1)
        img_path = match.group(2)
        start, end = match.span()
        
        if not img_path.startswith(('http://', 'https://', 'data:')):
            references.append({
                'alt': alt_text,
                'path': img_path,
                'start': start,
                'end': end
            })
    
    return references


def smart_match_images(image_references: List[Dict], extracted_images: List[Dict], md_images: List[Dict]) -> Dict:
    """
    智能匹配图片引用
    返回：{ md_image_index: image_info }
    """
    matches = {}
    
    placeholder_map = {}
    for ref in image_references:
        placeholder_name = ref.get('placeholder_name', '')
        placeholder_map[placeholder_name.lower()] = ref.get('image_info')
    
    for i, md_img in enumerate(md_images):
        md_path = md_img['path']
        md_filename = os.path.basename(md_path).lower()
        
        if md_filename in placeholder_map:
            matches[i] = placeholder_map[md_filename]
            continue
        
        for ref in image_references:
            img_info = ref.get('image_info', {})
            saved_name = img_info.get('saved_name', '').lower()
            original_name = img_info.get('original_name', '').lower()
            
            if md_filename == saved_name or md_filename == original_name:
                matches[i] = img_info
                break
    
    for i, md_img in enumerate(md_images):
        if i not in matches and i < len(extracted_images):
            matches[i] = extracted_images[i]
    
    return matches


def load_history() -> List[Dict]:
    if not HISTORY_FILE.exists():
        return []
    try:
        with open(HISTORY_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (json.JSONDecodeError, IOError):
        return []


def save_history(history: List[Dict]):
    try:
        with open(HISTORY_FILE, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
    except IOError:
        pass


def add_to_history(progress: ConversionProgress):
    history = load_history()
    
    history_item = {
        'task_id': progress.task_id,
        'filename': progress.filename,
        'status': progress.status,
        'created_at': progress.created_at.isoformat(),
        'completed_at': progress.completed_at.isoformat() if progress.completed_at else None,
        'output_format': progress.options.output_format,
        'images_extracted': progress.options.extract_images,
        'images': progress.images,
        'file_size': progress.original_file_size,
        'result_file': progress.result
    }
    
    history.insert(0, history_item)
    
    if len(history) > 100:
        history = history[:100]
    
    save_history(history)


def extract_and_save_images(
    markdown: str, 
    task_id: str, 
    file_path: Optional[Path] = None,
    options: Optional[ConversionOptions] = None
) -> Tuple[str, List[str]]:
    """
    提取并保存图片，确保路径完全匹配
    """
    all_extracted_images = []
    result = markdown
    
    task_images_dir = IMAGES_DIR / task_id
    task_images_dir.mkdir(exist_ok=True, parents=True)
    
    if options and options.extract_images:
        image_references = []
        extracted_from_file = []
        
        if file_path:
            file_ext = file_path.suffix.lower()
            
            if file_ext == '.pptx':
                extracted_from_file, image_references = extract_all_images_from_pptx(
                    file_path, task_images_dir, task_id
                )
                all_extracted_images.extend(extracted_from_file)
            
            elif file_ext == '.docx':
                extracted_from_file, image_references = extract_all_images_from_docx(
                    file_path, task_images_dir, task_id
                )
                all_extracted_images.extend(extracted_from_file)
        
        result, base64_images = extract_base64_images(result, task_images_dir, task_id)
        all_extracted_images.extend(base64_images)
        
        md_image_refs = find_all_image_references(result)
        
        if md_image_refs:
            matches = smart_match_images(image_references, all_extracted_images, md_image_refs)
            
            for i, md_img in enumerate(md_image_refs):
                if i in matches:
                    img_info = matches[i]
                    api_path = img_info.get('api_path', '')
                    
                    replacement = f"![{md_img['alt']}]({api_path})"
                    
                    offset = 0
                    for j in range(i):
                        if j in matches:
                            old_len = md_image_refs[j]['end'] - md_image_refs[j]['start']
                            new_len = len(f"![{md_image_refs[j]['alt']}]({matches[j].get('api_path', '')})")
                            offset += new_len - old_len
                    
                    start = md_img['start'] + offset
                    end = md_img['end'] + offset
                    
                    result = result[:start] + replacement + result[end:]
    
    img_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
    matches = list(re.finditer(img_pattern, result))
    
    for match in reversed(matches):
        alt_text = match.group(1)
        img_path = match.group(2)
        start, end = match.span()
        
        if img_path.startswith(('http://', 'https://', 'data:', '/api/images/')):
            continue
        
        original_img_path = Path(img_path)
        
        if original_img_path.is_absolute() and original_img_path.exists():
            new_img_name = f"{uuid.uuid4()}{original_img_path.suffix or '.png'}"
            new_img_path = task_images_dir / new_img_name
            
            try:
                shutil.copy2(original_img_path, new_img_path)
                
                api_url = f"/api/images/{task_id}/{new_img_name}"
                replacement = f"![{alt_text}]({api_url})"
                result = result[:start] + replacement + result[end:]
                
                all_extracted_images.append({
                    'saved_name': new_img_name,
                    'relative_path': f"images/{new_img_name}",
                    'api_path': api_url
                })
            except Exception:
                continue
        
        elif not original_img_path.is_absolute():
            found = False
            for img_info in all_extracted_images:
                if img_info.get('saved_name') == os.path.basename(img_path):
                    found = True
                    break
            
            if not found and file_path:
                parent_dir = file_path.parent
                possible_paths = [
                    parent_dir / img_path,
                    parent_dir / 'images' / os.path.basename(img_path),
                ]
                
                for possible_path in possible_paths:
                    if possible_path.exists():
                        new_img_name = f"{uuid.uuid4()}{possible_path.suffix or '.png'}"
                        new_img_path = task_images_dir / new_img_name
                        
                        try:
                            shutil.copy2(possible_path, new_img_path)
                            
                            api_url = f"/api/images/{task_id}/{new_img_name}"
                            replacement = f"![{alt_text}]({api_url})"
                            result = result[:start] + replacement + result[end:]
                            
                            all_extracted_images.append({
                                'saved_name': new_img_name,
                                'relative_path': f"images/{new_img_name}",
                                'api_path': api_url
                            })
                            found = True
                            break
                        except Exception:
                            continue
    
    relative_paths = [img.get('relative_path', '') for img in all_extracted_images if img.get('relative_path')]
    
    return result, relative_paths


async def convert_file_task(task_id: str, file_path: Path, filename: str, options: ConversionOptions):
    progress = conversion_tasks[task_id]
    
    try:
        progress.status = "converting"
        progress.progress = 10
        
        await asyncio.sleep(0.1)
        
        progress.progress = 30
        result = markitdown_instance.convert(str(file_path))
        
        progress.progress = 60
        content = result.text_content
        
        if options.optimize_tables:
            content = optimize_tables_format(content)
        
        if options.optimize_code_blocks:
            content = optimize_code_blocks_format(content)
        
        if options.output_format == "github":
            content = format_to_github_markdown(content)
        
        progress.progress = 70
        
        if options.extract_images:
            content, extracted_images = extract_and_save_images(
                content, 
                task_id, 
                file_path,
                options
            )
            progress.images = extracted_images
        
        progress.progress = 80
        
        result_filename = f"{task_id}.md"
        result_path = RESULT_DIR / result_filename
        
        with open(result_path, "w", encoding="utf-8") as f:
            f.write(content)
        
        progress.result = result_filename
        progress.status = "completed"
        progress.progress = 100
        progress.completed_at = datetime.now()
        
        add_to_history(progress)
        
    except UnsupportedFormatException as e:
        progress.status = "error"
        progress.error_message = f"Unsupported file format: {str(e)}. {get_supported_formats_message()}"
        progress.progress = 0
        add_to_history(progress)
        
    except FileConversionException as e:
        progress.status = "error"
        progress.error_message = f"File conversion failed: {str(e)}"
        progress.progress = 0
        add_to_history(progress)
        
    except MarkItDownException as e:
        progress.status = "error"
        progress.error_message = f"Conversion error: {str(e)}"
        progress.progress = 0
        add_to_history(progress)
        
    except Exception as e:
        progress.status = "error"
        progress.error_message = f"Unexpected error: {str(e)}"
        progress.progress = 0
        add_to_history(progress)
        
    finally:
        try:
            if file_path.exists():
                file_path.unlink()
        except Exception:
            pass


@app.get("/", response_class=HTMLResponse)
async def index():
    static_dir = Path(__file__).parent / "static"
    index_path = static_dir / "index.html"
    
    if index_path.exists():
        with open(index_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    
    return HTMLResponse(content="<h1>MarkItDown Web UI</h1><p>Static files not found.</p>")


@app.post("/api/upload", response_model=ConversionResponse)
async def upload_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    output_format: str = Form(default="standard"),
    extract_images: bool = Form(default=False),
    optimize_tables: bool = Form(default=False),
    optimize_code_blocks: bool = Form(default=False)
):
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")
    
    if not is_supported_format(file.filename):
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file format: {Path(file.filename).suffix}. {get_supported_formats_message()}"
        )
    
    task_id = str(uuid.uuid4())
    
    file_ext = Path(file.filename).suffix
    saved_filename = f"{task_id}{file_ext}"
    file_path = UPLOAD_DIR / saved_filename
    
    content = await file.read()
    file_size = len(content)
    with open(file_path, "wb") as f:
        f.write(content)
    
    options = ConversionOptions(
        output_format=output_format,
        extract_images=extract_images,
        optimize_tables=optimize_tables,
        optimize_code_blocks=optimize_code_blocks
    )
    
    progress = ConversionProgress(
        task_id=task_id,
        filename=file.filename,
        status="pending",
        progress=0,
        options=options,
        original_file_size=file_size,
    )
    conversion_tasks[task_id] = progress
    
    background_tasks.add_task(convert_file_task, task_id, file_path, file.filename, options)
    
    return ConversionResponse(
        task_id=task_id,
        filename=file.filename,
        status="pending",
        progress=0,
    )


@app.post("/api/upload/batch", response_model=BatchConversionResponse)
async def upload_batch_files(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
    output_format: str = Form(default="standard"),
    extract_images: bool = Form(default=False),
    optimize_tables: bool = Form(default=False),
    optimize_code_blocks: bool = Form(default=False)
):
    responses = []
    options = ConversionOptions(
        output_format=output_format,
        extract_images=extract_images,
        optimize_tables=optimize_tables,
        optimize_code_blocks=optimize_code_blocks
    )
    
    for file in files:
        if not file.filename:
            continue
        
        if not is_supported_format(file.filename):
            continue
        
        task_id = str(uuid.uuid4())
        
        file_ext = Path(file.filename).suffix
        saved_filename = f"{task_id}{file_ext}"
        file_path = UPLOAD_DIR / saved_filename
        
        content = await file.read()
        file_size = len(content)
        with open(file_path, "wb") as f:
            f.write(content)
        
        progress = ConversionProgress(
            task_id=task_id,
            filename=file.filename,
            status="pending",
            progress=0,
            options=options,
            original_file_size=file_size,
        )
        conversion_tasks[task_id] = progress
        
        background_tasks.add_task(convert_file_task, task_id, file_path, file.filename, options)
        
        responses.append(ConversionResponse(
            task_id=task_id,
            filename=file.filename,
            status="pending",
            progress=0,
        ))
    
    return BatchConversionResponse(tasks=responses)


@app.get("/api/progress/{task_id}", response_model=ProgressResponse)
async def get_progress(task_id: str):
    if task_id not in conversion_tasks:
        raise HTTPException(status_code=404, detail="Task not found")
    
    progress = conversion_tasks[task_id]
    
    return ProgressResponse(
        task_id=task_id,
        filename=progress.filename,
        status=progress.status,
        progress=progress.progress,
        error_message=progress.error_message,
        images=progress.images if progress.images else None,
    )


def create_zip_package(task_id: str, progress: ConversionProgress) -> io.BytesIO:
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        result_path = RESULT_DIR / progress.result
        if result_path.exists():
            with open(result_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            md_content = re.sub(
                r'!\[([^\]]*)\]\(/api/images/[^/]+/([^)]+)\)',
                r'![\1](images/\2)',
                md_content
            )
            
            original_filename = Path(progress.filename).stem
            safe_filename = sanitize_filename(original_filename)
            md_filename = f"{safe_filename}.md"
            
            zip_file.writestr(md_filename, md_content.encode('utf-8'))
        
        task_images_dir = IMAGES_DIR / task_id
        if task_images_dir.exists():
            for img_file in task_images_dir.iterdir():
                if img_file.is_file():
                    img_filename = sanitize_filename(img_file.name)
                    arcname = f"images/{img_filename}"
                    zip_file.write(str(img_file), arcname)
    
    zip_buffer.seek(0)
    return zip_buffer


@app.get("/api/download/{task_id}")
async def download_result(task_id: str, format: str = "auto"):
    if task_id not in conversion_tasks:
        raise HTTPException(status_code=404, detail="Task not found")
    
    progress = conversion_tasks[task_id]
    
    if progress.status != "completed" or not progress.result:
        raise HTTPException(status_code=400, detail="Conversion not completed yet")
    
    result_path = RESULT_DIR / progress.result
    
    if not result_path.exists():
        raise HTTPException(status_code=404, detail="Result file not found")
    
    has_images = progress.images and len(progress.images) > 0
    
    if format == "zip" or (format == "auto" and has_images):
        zip_buffer = create_zip_package(task_id, progress)
        
        original_filename = Path(progress.filename).stem
        safe_filename = sanitize_filename(original_filename)
        download_filename = f"{safe_filename}.zip"
        
        return StreamingResponse(
            io.BytesIO(zip_buffer.getvalue()),
            media_type="application/zip",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{download_filename}"
            }
        )
    else:
        with open(result_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        md_content = re.sub(
            r'!\[([^\]]*)\]\(/api/images/[^/]+/([^)]+)\)',
            r'![\1](images/\2)',
            md_content
        )
        
        original_filename = Path(progress.filename).stem
        safe_filename = sanitize_filename(original_filename)
        download_filename = f"{safe_filename}.md"
        
        temp_md_path = RESULT_DIR / f"{task_id}_download.md"
        with open(temp_md_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        
        return FileResponse(
            path=str(temp_md_path),
            media_type="text/markdown",
            filename=download_filename,
        )


@app.get("/api/download/{task_id}/zip")
async def download_result_zip(task_id: str):
    return await download_result(task_id, format="zip")


@app.get("/api/result/{task_id}")
async def get_result_content(task_id: str):
    if task_id not in conversion_tasks:
        raise HTTPException(status_code=404, detail="Task not found")
    
    progress = conversion_tasks[task_id]
    
    if progress.status != "completed" or not progress.result:
        raise HTTPException(status_code=400, detail="Conversion not completed yet")
    
    result_path = RESULT_DIR / progress.result
    
    if not result_path.exists():
        raise HTTPException(status_code=404, detail="Result file not found")
    
    with open(result_path, "r", encoding="utf-8") as f:
        content = f.read()
    
    return JSONResponse(content={
        "task_id": task_id,
        "filename": progress.filename,
        "content": content,
    })


@app.get("/api/formats")
async def get_supported_formats():
    return JSONResponse(content={
        "supported_extensions": sorted(list(set(SUPPORTED_EXTENSIONS))),
        "message": get_supported_formats_message(),
    })


def cleanup_old_files():
    try:
        cutoff = datetime.now().timestamp() - 3600
        
        for file_path in UPLOAD_DIR.iterdir():
            if file_path.is_file():
                if file_path.stat().st_mtime < cutoff:
                    file_path.unlink()
        
        for file_path in RESULT_DIR.iterdir():
            if file_path.is_file():
                if file_path.stat().st_mtime < cutoff:
                    file_path.unlink()
        
        cutoff_time = datetime.now().timestamp() - 3600
        tasks_to_remove = []
        for task_id, progress in conversion_tasks.items():
            if progress.completed_at:
                if progress.completed_at.timestamp() < cutoff_time:
                    tasks_to_remove.append(task_id)
        
        for task_id in tasks_to_remove:
            del conversion_tasks[task_id]
            
    except Exception:
        pass


@app.get("/api/images/{task_id}/{filename}")
async def get_extracted_image(task_id: str, filename: str):
    decoded_filename = unquote(filename)
    sanitized_filename = sanitize_filename(decoded_filename)
    
    image_path = IMAGES_DIR / task_id / sanitized_filename
    
    if not image_path.exists():
        task_images_dir = IMAGES_DIR / task_id
        if task_images_dir.exists():
            for img_file in task_images_dir.iterdir():
                if img_file.is_file():
                    if img_file.name == sanitized_filename:
                        image_path = img_file
                        break
                    
                    if img_file.name.lower() == sanitized_filename.lower():
                        image_path = img_file
                        break
                    
                    decoded_file = unquote(img_file.name)
                    if decoded_file == decoded_filename:
                        image_path = img_file
                        break
    
    if not image_path.exists():
        raise HTTPException(status_code=404, detail="Image not found")
    
    content_type = "image/png"
    ext = image_path.suffix.lower()
    if ext in ('.jpg', '.jpeg'):
        content_type = "image/jpeg"
    elif ext == '.gif':
        content_type = "image/gif"
    elif ext == '.bmp':
        content_type = "image/bmp"
    elif ext == '.webp':
        content_type = "image/webp"
    elif ext == '.svg':
        content_type = "image/svg+xml"
    elif ext == '.ico':
        content_type = "image/x-icon"
    
    return FileResponse(
        path=str(image_path),
        media_type=content_type
    )


@app.get("/api/history", response_model=List[HistoryItem])
async def get_conversion_history(limit: int = 50):
    history = load_history()
    limited_history = history[:limit]
    
    return [
        HistoryItem(
            task_id=item['task_id'],
            filename=item['filename'],
            status=item['status'],
            created_at=item['created_at'],
            completed_at=item.get('completed_at'),
            output_format=item.get('output_format', 'standard'),
            images_extracted=item.get('images_extracted', False),
            file_size=item.get('file_size', 0)
        )
        for item in limited_history
    ]


@app.delete("/api/history/{task_id}")
async def delete_history_item(task_id: str):
    history = load_history()
    item_to_remove = None
    
    for i, item in enumerate(history):
        if item['task_id'] == task_id:
            item_to_remove = item
            break
    
    if item_to_remove is None:
        raise HTTPException(status_code=404, detail="History item not found")
    
    history = [item for item in history if item['task_id'] != task_id]
    save_history(history)
    
    if item_to_remove.get('result_file'):
        result_path = RESULT_DIR / item_to_remove['result_file']
        if result_path.exists():
            result_path.unlink()
    
    task_images_dir = IMAGES_DIR / task_id
    if task_images_dir.exists():
        shutil.rmtree(task_images_dir)
    
    return JSONResponse(content={
        "success": True,
        "message": f"History item {task_id} deleted"
    })


@app.delete("/api/history")
async def clear_all_history():
    history = load_history()
    
    for item in history:
        if item.get('result_file'):
            result_path = RESULT_DIR / item['result_file']
            if result_path.exists():
                try:
                    result_path.unlink()
                except Exception:
                    pass
        
        task_images_dir = IMAGES_DIR / item['task_id']
        if task_images_dir.exists():
            try:
                shutil.rmtree(task_images_dir)
            except Exception:
                pass
    
    if HISTORY_FILE.exists():
        HISTORY_FILE.unlink()
    
    return JSONResponse(content={
        "success": True,
        "message": "All history cleared",
        "items_cleared": len(history)
    })


@app.post("/api/cache/clear", response_model=ClearCacheResponse)
async def clear_cache():
    uploads_cleared = 0
    results_cleared = 0
    images_cleared = 0
    
    try:
        if UPLOAD_DIR.exists():
            for file_path in UPLOAD_DIR.iterdir():
                if file_path.is_file():
                    file_path.unlink()
                    uploads_cleared += 1
                elif file_path.is_dir():
                    shutil.rmtree(file_path)
                    uploads_cleared += 1
    except Exception:
        pass
    
    try:
        if RESULT_DIR.exists():
            for file_path in RESULT_DIR.iterdir():
                if file_path.is_file() and file_path.suffix == '.md':
                    file_path.unlink()
                    results_cleared += 1
    except Exception:
        pass
    
    try:
        if IMAGES_DIR.exists():
            for dir_path in IMAGES_DIR.iterdir():
                if dir_path.is_dir():
                    shutil.rmtree(dir_path)
                    images_cleared += 1
    except Exception:
        pass
    
    conversion_tasks.clear()
    
    return ClearCacheResponse(
        success=True,
        message="Cache cleared successfully",
        uploads_cleared=uploads_cleared,
        results_cleared=results_cleared,
        images_cleared=images_cleared
    )


@app.get("/api/cache/stats")
async def get_cache_stats():
    upload_count = 0
    upload_size = 0
    result_count = 0
    result_size = 0
    image_count = 0
    image_size = 0
    
    try:
        if UPLOAD_DIR.exists():
            for file_path in UPLOAD_DIR.iterdir():
                if file_path.is_file():
                    upload_count += 1
                    upload_size += file_path.stat().st_size
    except Exception:
        pass
    
    try:
        if RESULT_DIR.exists():
            for file_path in RESULT_DIR.iterdir():
                if file_path.is_file() and file_path.suffix == '.md':
                    result_count += 1
                    result_size += file_path.stat().st_size
    except Exception:
        pass
    
    try:
        if IMAGES_DIR.exists():
            for task_dir in IMAGES_DIR.iterdir():
                if task_dir.is_dir():
                    for img_path in task_dir.iterdir():
                        if img_path.is_file():
                            image_count += 1
                            image_size += img_path.stat().st_size
    except Exception:
        pass
    
    def format_size(size: int) -> str:
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size < 1024:
                return f"{size:.1f} {unit}"
            size /= 1024
        return f"{size:.1f} TB"
    
    return JSONResponse(content={
        "uploads": {
            "count": upload_count,
            "size": upload_size,
            "size_formatted": format_size(upload_size)
        },
        "results": {
            "count": result_count,
            "size": result_size,
            "size_formatted": format_size(result_size)
        },
        "images": {
            "count": image_count,
            "size": image_size,
            "size_formatted": format_size(image_size)
        },
        "history_count": len(load_history()),
        "active_tasks": len(conversion_tasks)
    })


static_dir = Path(__file__).parent / "static"
if static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")

images_abs_dir = Path(__file__).parent.parent.parent.parent.parent / IMAGES_DIR
if not images_abs_dir.exists():
    images_abs_dir = IMAGES_DIR
    
if images_abs_dir.exists():
    app.mount("/images", StaticFiles(directory=str(images_abs_dir)), name="extracted_images")
