"""
Enhanced PPTX Converter with improved OCR support.
Already has LLM-based image description, this enhances it with traditional OCR fallback.
"""

import io
import os
import sys
import tempfile
from typing import Any, BinaryIO, Optional

from typing import BinaryIO, Any, Optional

from markitdown.converters import HtmlConverter
from markitdown import DocumentConverter, DocumentConverterResult, StreamInfo
from markitdown._exceptions import (
    MissingDependencyException,
    MISSING_DEPENDENCY_MESSAGE,
)
from ._ocr_service import LLMVisionOCRService, format_image_reference

_dependency_exc_info = None
try:
    import pptx
except ImportError:
    _dependency_exc_info = sys.exc_info()


class PptxConverterWithOCR(DocumentConverter):
    """Enhanced PPTX Converter with OCR fallback."""

    def __init__(self, ocr_service: Optional[LLMVisionOCRService] = None):
        super().__init__()
        self._html_converter = HtmlConverter()
        self.ocr_service = ocr_service

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension == ".pptx":
            return True

        if mimetype.startswith(
            "application/vnd.openxmlformats-officedocument.presentationml"
        ):
            return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> DocumentConverterResult:
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[1].with_traceback(
                _dependency_exc_info[2]
            )  # type: ignore[union-attr]

        # Get OCR service (from kwargs or instance)
        ocr_service: Optional[LLMVisionOCRService] = (
            kwargs.get("ocr_service") or self.ocr_service
        )
        llm_client = kwargs.get("llm_client")

        # --- extract_only mode: skip OCR, emit image file references ---
        if kwargs.get("extract_only", False):
            image_output_dir = kwargs.get("image_output_dir") or tempfile.mkdtemp(
                prefix="markitdown_ocr_"
            )
            os.makedirs(image_output_dir, exist_ok=True)
            # Pop keys already consumed as positional args to avoid duplicate keyword
            _eo_kwargs = {k: v for k, v in kwargs.items() if k not in ("image_output_dir",)}
            return self._convert_extract_only(file_stream, image_output_dir, **_eo_kwargs)

        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0

        for slide in presentation.slides:
            slide_num += 1
            md_content += f"\\n\\n<!-- Slide number: {slide_num} -->\\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content

                # Pictures
                if self._is_picture(shape):
                    # Get image data
                    image_stream = io.BytesIO(shape.image.blob)

                    # Try LLM description first if available
                    llm_description = ""
                    if llm_client and kwargs.get("llm_model"):
                        try:
                            from ._llm_caption import llm_caption

                            image_filename = shape.image.filename
                            image_extension = None
                            if image_filename:
                                import os

                                image_extension = os.path.splitext(image_filename)[1]

                            image_stream_info = StreamInfo(
                                mimetype=shape.image.content_type,
                                extension=image_extension,
                                filename=image_filename,
                            )

                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=kwargs.get("llm_model"),
                                prompt=kwargs.get("llm_prompt"),
                            )
                        except Exception:
                            pass

                    # Try OCR if LLM failed or not available
                    ocr_text = ""
                    if not llm_description and ocr_service:
                        try:
                            image_stream.seek(0)
                            ocr_result = ocr_service.extract_text(image_stream)
                            if ocr_result.text.strip():
                                ocr_text = ocr_result.text.strip()
                        except Exception:
                            pass

                    # Format extracted content using unified OCR block format
                    content = (llm_description or ocr_text or "").strip()
                    if content:
                        md_content += f"\n*[Image OCR]\n{content}\n[End OCR]*\n"

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\\n"
                    else:
                        md_content += shape.text + "\\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\\n\\n### Notes:\\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _convert_extract_only(
        self, file_stream: BinaryIO, image_output_dir: str, **kwargs: Any
    ) -> DocumentConverterResult:
        """
        Extract-only mode: extract text and save embedded images to disk.
        No OCR or LLM description is performed; images are referenced via file paths.
        """
        from PIL import Image

        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0
        global_img_idx = 0

        for slide in presentation.slides:
            slide_num += 1
            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            def get_shape_content_extract_only(shape, **kw):
                nonlocal md_content, global_img_idx

                # Pictures
                if self._is_picture(shape):
                    try:
                        image_bytes = shape.image.blob

                        ext = "png"
                        width, height = None, None
                        try:
                            pil_img = Image.open(io.BytesIO(image_bytes))
                            fmt = pil_img.format
                            if fmt:
                                ext = fmt.lower()
                                if ext == "jpeg":
                                    ext = "jpg"
                            width, height = pil_img.size
                        except Exception:
                            pass

                        filename = f"slide_{slide_num}_{global_img_idx}.{ext}"
                        filepath = os.path.join(image_output_dir, filename)
                        with open(filepath, "wb") as f:
                            f.write(image_bytes)

                        img_ref = format_image_reference(
                            filepath,
                            width=width,
                            height=height,
                            size_bytes=len(image_bytes),
                        )
                        md_content += f"\n{img_ref}\n"
                        global_img_idx += 1
                    except Exception:
                        pass

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kw)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content_extract_only(subshape, **kw)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content_extract_only(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        import html

        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\\n\\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\\n\\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            if "unsupported plot type" in str(e):
                return "\\n\\n[unsupported chart]\\n\\n"
        except Exception:
            return "\\n\\n[unsupported chart]\\n\\n"
