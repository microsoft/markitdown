# type: ignore
import base64
import binascii
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import zipfile
import importlib
import sys
from importlib.metadata import entry_points
from xml.dom import minidom
from typing import Any, Dict, List, Optional, Union
from pathlib import Path
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse
from warnings import warn, resetwarnings, catch_warnings

import mammoth
import markdownify
import olefile
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx

# File-format detection
import puremagic
import requests
from bs4 import BeautifulSoup
from charset_normalizer import from_path

# Azure imports
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.ai.documentintelligence.models import (
    AnalyzeDocumentRequest,
    AnalyzeResult,
    DocumentAnalysisFeature,
)
from azure.identity import DefaultAzureCredential

from .converters import (
    DocumentConverter,
    DocumentConverterResult,
    PlainTextConverter,
    HtmlConverter,
    RssConverter,
    WikipediaConverter,
    YouTubeConverter,
)
from .converters._markdownify import _CustomMarkdownify

from ._exceptions import (
    MarkItDownException,
    ConverterPrerequisiteException,
    FileConversionException,
    UnsupportedFormatException,
)

# TODO: currently, there is a bug in the document intelligence SDK with importing the "ContentFormat" enum.
# This constant is a temporary fix until the bug is resolved.
CONTENT_FORMAT = "markdown"

# Override mimetype for csv to fix issue on windows
mimetypes.add_type("text/csv", ".csv")

PRIORITY_SPECIFIC_FILE_FORMAT = 0.0
PRIORITY_GENERIC_FILE_FORMAT = -10.0

# Optional Transcription support
IS_AUDIO_TRANSCRIPTION_CAPABLE = False
try:
    # Using warnings' catch_warnings to catch
    # pydub's warning of ffmpeg or avconv missing
    with catch_warnings(record=True) as w:
        import pydub

        if w:
            raise ModuleNotFoundError
    import speech_recognition as sr

    IS_AUDIO_TRANSCRIPTION_CAPABLE = True
except ModuleNotFoundError:
    pass
finally:
    resetwarnings()


class IpynbConverter(DocumentConverter):
    """Converts Jupyter Notebook (.ipynb) files to Markdown."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not ipynb
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".ipynb":
            return None

        # Parse and convert the notebook
        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            notebook_content = json.load(fh)
            result = self._convert(notebook_content)

        return result

    def _convert(self, notebook_content: dict) -> Union[None, DocumentConverterResult]:
        """Helper function that converts notebook JSON content to Markdown."""
        try:
            md_output = []
            title = None

            for cell in notebook_content.get("cells", []):
                cell_type = cell.get("cell_type", "")
                source_lines = cell.get("source", [])

                if cell_type == "markdown":
                    md_output.append("".join(source_lines))

                    # Extract the first # heading as title if not already found
                    if title is None:
                        for line in source_lines:
                            if line.startswith("# "):
                                title = line.lstrip("# ").strip()
                                break

                elif cell_type == "code":
                    # Code cells are wrapped in Markdown code blocks
                    md_output.append(f"```python\n{''.join(source_lines)}\n```")
                elif cell_type == "raw":
                    md_output.append(f"```\n{''.join(source_lines)}\n```")

            md_text = "\n\n".join(md_output)

            # Check for title in notebook metadata
            title = notebook_content.get("metadata", {}).get("title", title)

            return DocumentConverterResult(
                title=title,
                text_content=md_text,
            )

        except Exception as e:
            raise FileConversionException(
                f"Error converting .ipynb file: {str(e)}"
            ) from e


class BingSerpConverter(DocumentConverter):
    """
    Handle Bing results pages (only the organic search results).
    NOTE: It is better to use the Bing API
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a Bing SERP
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https://www\.bing\.com/search\?q=", url):
            return None

        # Parse the query parameters
        parsed_params = parse_qs(urlparse(url).query)
        query = parsed_params.get("q", [""])[0]

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Clean up some formatting
        for tptt in soup.find_all(class_="tptt"):
            if hasattr(tptt, "string") and tptt.string:
                tptt.string += " "
        for slug in soup.find_all(class_="algoSlug_icon"):
            slug.extract()

        # Parse the algorithmic results
        _markdownify = _CustomMarkdownify()
        results = list()
        for result in soup.find_all(class_="b_algo"):
            # Rewrite redirect urls
            for a in result.find_all("a", href=True):
                parsed_href = urlparse(a["href"])
                qs = parse_qs(parsed_href.query)

                # The destination is contained in the u parameter,
                # but appears to be base64 encoded, with some prefix
                if "u" in qs:
                    u = (
                        qs["u"][0][2:].strip() + "=="
                    )  # Python 3 doesn't care about extra padding

                    try:
                        # RFC 4648 / Base64URL" variant, which uses "-" and "_"
                        a["href"] = base64.b64decode(u, altchars="-_").decode("utf-8")
                    except UnicodeDecodeError:
                        pass
                    except binascii.Error:
                        pass

            # Convert to markdown
            md_result = _markdownify.convert_soup(result).strip()
            lines = [line.strip() for line in re.split(r"\n+", md_result)]
            results.append("\n".join([line for line in lines if len(line) > 0]))

        webpage_text = (
            f"## A Bing search for '{query}' found the following results:\n\n"
            + "\n\n".join(results)
        )

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            style_map = kwargs.get("style_map", None)

            result = mammoth.convert_to_html(docx_file, style_map=style_map)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None, engine="openpyxl")
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class XlsConverter(HtmlConverter):
    """
    Converts XLS files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLS
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xls":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None, engine="xlrd")
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def _get_llm_description(
        self, llm_client, llm_model, image_blob, content_type, prompt=None
    ):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed alt text for this image with less than 50 words."

        image_base64 = base64.b64encode(image_blob).decode("utf-8")
        data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                    {"type": "text", "text": prompt},
                ],
            }
        ]

        response = llm_client.chat.completions.create(
            model=llm_model, messages=messages
        )
        return response.choices[0].message.content

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    llm_description = None
                    alt_text = None

                    llm_client = kwargs.get("llm_client")
                    llm_model = kwargs.get("llm_model")
                    if llm_client is not None and llm_model is not None:
                        try:
                            llm_description = self._get_llm_description(
                                llm_client,
                                llm_model,
                                shape.image.blob,
                                shape.image.content_type,
                            )
                        except Exception:
                            # Unable to describe with LLM
                            pass

                    if not llm_description:
                        try:
                            alt_text = shape._element._nvXxPr.cNvPr.attrib.get(
                                "descr", ""
                            )
                        except Exception:
                            # Unable to get alt text
                            pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += (
                        "\n!["
                        + (llm_description or alt_text or shape.name)
                        + "]("
                        + filename
                        + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_chart_to_markdown(self, chart):
        md = "\n\n### Chart"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"
        data = []
        category_names = [c.label for c in chart.plots[0].categories]
        series_names = [s.name for s in chart.series]
        data.append(["Category"] + series_names)

        for idx, category in enumerate(category_names):
            row = [category]
            for series in chart.series:
                row.append(series.values[idx])
            data.append(row)

        markdown_table = []
        for row in data:
            markdown_table.append("| " + " | ".join(map(str, row)) + " |")
        header = markdown_table[0]
        separator = "|" + "|".join(["---"] * len(data[0])) + "|"
        return md + "\n".join([header, separator] + markdown_table[1:])


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path, exiftool_path=None):
        if not exiftool_path:
            which_exiftool = shutil.which("exiftool")
            if which_exiftool:
                warn(
                    f"""Implicit discovery of 'exiftool' is disabled. If you would like to continue to use exiftool in MarkItDown, please set the exiftool_path parameter in the MarkItDown consructor. E.g., 

    md = MarkItDown(exiftool_path="{which_exiftool}")

This warning will be removed in future releases.
""",
                    DeprecationWarning,
                )

            return None
        else:
            try:
                result = subprocess.run(
                    [exiftool_path, "-json", local_path], capture_output=True, text=True
                ).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a WAV
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path, kwargs.get("exiftool_path"))
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            try:
                transcript = self._transcribe_audio(local_path)
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += (
                    "\n\n### Audio Transcript:\nError. Could not transcribe this audio."
                )

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".mp3":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path, kwargs.get("exiftool_path"))
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            handle, temp_path = tempfile.mkstemp(suffix=".wav")
            os.close(handle)
            try:
                sound = pydub.AudioSegment.from_mp3(local_path)
                sound.export(temp_path, format="wav")

                _args = dict()
                _args.update(kwargs)
                _args["file_extension"] = ".wav"

                try:
                    transcript = super()._transcribe_audio(temp_path).strip()
                    md_content += "\n\n### Audio Transcript:\n" + (
                        "[No speech detected]" if transcript == "" else transcript
                    )
                except Exception:
                    md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

            finally:
                os.unlink(temp_path)

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an llm_client is configured).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not an image
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path, kwargs.get("exiftool_path"))
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        llm_client = kwargs.get("llm_client")
        llm_model = kwargs.get("llm_model")
        if llm_client is not None and llm_model is not None:
            md_content += (
                "\n# Description:\n"
                + self._get_llm_description(
                    local_path,
                    extension,
                    llm_client,
                    llm_model,
                    prompt=kwargs.get("llm_prompt"),
                ).strip()
                + "\n"
            )

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _get_llm_description(self, local_path, extension, client, model, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                ],
            }
        ]

        response = client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content


class OutlookMsgConverter(DocumentConverter):
    """Converts Outlook .msg files to markdown by extracting email metadata and content.

    Uses the olefile package to parse the .msg file structure and extract:
    - Email headers (From, To, Subject)
    - Email body content
    """

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not a MSG file
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".msg":
            return None

        try:
            msg = olefile.OleFileIO(local_path)
            # Extract email metadata
            md_content = "# Email Message\n\n"

            # Get headers
            headers = {
                "From": self._get_stream_data(msg, "__substg1.0_0C1F001F"),
                "To": self._get_stream_data(msg, "__substg1.0_0E04001F"),
                "Subject": self._get_stream_data(msg, "__substg1.0_0037001F"),
            }

            # Add headers to markdown
            for key, value in headers.items():
                if value:
                    md_content += f"**{key}:** {value}\n"

            md_content += "\n## Content\n\n"

            # Get email body
            body = self._get_stream_data(msg, "__substg1.0_1000001F")
            if body:
                md_content += body

            msg.close()

            return DocumentConverterResult(
                title=headers.get("Subject"), text_content=md_content.strip()
            )

        except Exception as e:
            raise FileConversionException(
                f"Could not convert MSG file '{local_path}': {str(e)}"
            )

    def _get_stream_data(
        self, msg: olefile.OleFileIO, stream_path: str
    ) -> Union[str, None]:
        """Helper to safely extract and decode stream data from the MSG file."""
        try:
            if msg.exists(stream_path):
                data = msg.openstream(stream_path).read()
                # Try UTF-16 first (common for .msg files)
                try:
                    return data.decode("utf-16-le").strip()
                except UnicodeDecodeError:
                    # Fall back to UTF-8
                    try:
                        return data.decode("utf-8").strip()
                    except UnicodeDecodeError:
                        # Last resort - ignore errors
                        return data.decode("utf-8", errors="ignore").strip()
        except Exception:
            pass
        return None


class ZipConverter(DocumentConverter):
    """Converts ZIP files to markdown by extracting and converting all contained files.

    The converter extracts the ZIP contents to a temporary directory, processes each file
    using appropriate converters based on file extensions, and then combines the results
    into a single markdown document. The temporary directory is cleaned up after processing.

    Example output format:
    ```markdown
    Content from the zip file `example.zip`:

    ## File: docs/readme.txt

    This is the content of readme.txt
    Multiple lines are preserved

    ## File: images/example.jpg

    ImageSize: 1920x1080
    DateTimeOriginal: 2024-02-15 14:30:00
    Description: A beautiful landscape photo

    ## File: data/report.xlsx

    ## Sheet1
    | Column1 | Column2 | Column3 |
    |---------|---------|---------|
    | data1   | data2   | data3   |
    | data4   | data5   | data6   |
    ```

    Key features:
    - Maintains original file structure in headings
    - Processes nested files recursively
    - Uses appropriate converters for each file type
    - Preserves formatting of converted content
    - Cleans up temporary files after processing
    """

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not a ZIP
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        # Get parent converters list if available
        parent_converters = kwargs.get("_parent_converters", [])
        if not parent_converters:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] No converters available to process zip contents from: {local_path}",
            )

        extracted_zip_folder_name = (
            f"extracted_{os.path.basename(local_path).replace('.zip', '_zip')}"
        )
        extraction_dir = os.path.normpath(
            os.path.join(os.path.dirname(local_path), extracted_zip_folder_name)
        )
        md_content = f"Content from the zip file `{os.path.basename(local_path)}`:\n\n"

        try:
            # Extract the zip file safely
            with zipfile.ZipFile(local_path, "r") as zipObj:
                # Safeguard against path traversal
                for member in zipObj.namelist():
                    member_path = os.path.normpath(os.path.join(extraction_dir, member))
                    if (
                        not os.path.commonprefix([extraction_dir, member_path])
                        == extraction_dir
                    ):
                        raise ValueError(
                            f"Path traversal detected in zip file: {member}"
                        )

                # Extract all files safely
                zipObj.extractall(path=extraction_dir)

            # Process each extracted file
            for root, dirs, files in os.walk(extraction_dir):
                for name in files:
                    file_path = os.path.join(root, name)
                    relative_path = os.path.relpath(file_path, extraction_dir)

                    # Get file extension
                    _, file_extension = os.path.splitext(name)

                    # Update kwargs for the file
                    file_kwargs = kwargs.copy()
                    file_kwargs["file_extension"] = file_extension
                    file_kwargs["_parent_converters"] = parent_converters

                    # Try converting the file using available converters
                    for converter in parent_converters:
                        # Skip the zip converter to avoid infinite recursion
                        if isinstance(converter, ZipConverter):
                            continue

                        result = converter.convert(file_path, **file_kwargs)
                        if result is not None:
                            md_content += f"\n## File: {relative_path}\n\n"
                            md_content += result.text_content + "\n\n"
                            break

            # Clean up extracted files if specified
            if kwargs.get("cleanup_extracted", True):
                shutil.rmtree(extraction_dir)

            return DocumentConverterResult(title=None, text_content=md_content.strip())

        except zipfile.BadZipFile:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] Invalid or corrupted zip file: {local_path}",
            )
        except ValueError as ve:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] Security error in zip file {local_path}: {str(ve)}",
            )
        except Exception as e:
            return DocumentConverterResult(
                title=None,
                text_content=f"[ERROR] Failed to process zip file {local_path}: {str(e)}",
            )


class DocumentIntelligenceConverter(DocumentConverter):
    """Specialized DocumentConverter that uses Document Intelligence to extract text from documents."""

    def __init__(
        self,
        endpoint: str,
        api_version: str = "2024-07-31-preview",
    ):
        self.endpoint = endpoint
        self.api_version = api_version
        self.doc_intel_client = DocumentIntelligenceClient(
            endpoint=self.endpoint,
            api_version=self.api_version,
            credential=DefaultAzureCredential(),
        )

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if extension is not supported by Document Intelligence
        extension = kwargs.get("file_extension", "")
        docintel_extensions = [
            ".pdf",
            ".docx",
            ".xlsx",
            ".pptx",
            ".html",
            ".jpeg",
            ".jpg",
            ".png",
            ".bmp",
            ".tiff",
            ".heif",
        ]
        if extension.lower() not in docintel_extensions:
            return None

        # Get the bytestring for the local path
        with open(local_path, "rb") as f:
            file_bytes = f.read()

        # Certain document analysis features are not availiable for filetypes (.xlsx, .pptx, .html)
        if extension.lower() in [".xlsx", ".pptx", ".html"]:
            analysis_features = []
        else:
            analysis_features = [
                DocumentAnalysisFeature.FORMULAS,  # enable formula extraction
                DocumentAnalysisFeature.OCR_HIGH_RESOLUTION,  # enable high resolution OCR
                DocumentAnalysisFeature.STYLE_FONT,  # enable font style extraction
            ]

        # Extract the text using Azure Document Intelligence
        poller = self.doc_intel_client.begin_analyze_document(
            model_id="prebuilt-layout",
            body=AnalyzeDocumentRequest(bytes_source=file_bytes),
            features=analysis_features,
            output_content_format=CONTENT_FORMAT,  # TODO: replace with "ContentFormat.MARKDOWN" when the bug is fixed
        )
        result: AnalyzeResult = poller.result()

        # remove comments from the markdown content generated by Doc Intelligence and append to markdown string
        markdown_text = re.sub(r"<!--.*?-->", "", result.content, flags=re.DOTALL)
        return DocumentConverterResult(
            title=None,
            text_content=markdown_text,
        )


class MarkItDown:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        llm_client: Optional[Any] = None,
        llm_model: Optional[str] = None,
        style_map: Optional[str] = None,
        exiftool_path: Optional[str] = None,
        docintel_endpoint: Optional[str] = None,
        # Deprecated
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[str] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        if exiftool_path is None:
            exiftool_path = os.environ.get("EXIFTOOL_PATH")

        # Handle deprecation notices
        #############################
        if mlm_client is not None:
            if llm_client is None:
                warn(
                    "'mlm_client' is deprecated, and was renamed 'llm_client'.",
                    DeprecationWarning,
                )
                llm_client = mlm_client
                mlm_client = None
            else:
                raise ValueError(
                    "'mlm_client' is deprecated, and was renamed 'llm_client'. Do not use both at the same time. Just use 'llm_client' instead."
                )

        if mlm_model is not None:
            if llm_model is None:
                warn(
                    "'mlm_model' is deprecated, and was renamed 'llm_model'.",
                    DeprecationWarning,
                )
                llm_model = mlm_model
                mlm_model = None
            else:
                raise ValueError(
                    "'mlm_model' is deprecated, and was renamed 'llm_model'. Do not use both at the same time. Just use 'llm_model' instead."
                )
        #############################

        self._llm_client = llm_client
        self._llm_model = llm_model
        self._style_map = style_map
        self._exiftool_path = exiftool_path

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(RssConverter())
        self.register_page_converter(WikipediaConverter())

        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(BingSerpConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(XlsConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(IpynbConverter())
        self.register_page_converter(PdfConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(OutlookMsgConverter())

        #        print("Discovering plugins")
        #        for entry_point in entry_points(group="markitdown.converters"):
        #            args = {
        #                "required1": "Override1",
        #                "required2": "Override2",
        #                "required3": "Override3"
        #            }
        #
        #            #print(entry_point)
        #            plugin = entry_point.load()
        #            instance = plugin(**args)
        #            print(instance)

        #    try:
        #        ConverterClass = entry_point.load()
        #        self.register_page_converter(ConverterClass())
        #        print(f"✔ Registered converter: {entry_point.name}")
        #    except Exception as e:
        #        print(f" Failed to load {entry_point.name}: {e}")
        #        print("Done")

        # Register Document Intelligence converter at the top of the stack if endpoint is provided
        if docintel_endpoint is not None:
            self.register_page_converter(
                DocumentIntelligenceConverter(endpoint=docintel_endpoint)
            )

    def convert(
        self, source: Union[str, requests.Response, Path], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path either as string pathlib path object or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if (
                source.startswith("http://")
                or source.startswith("https://")
                or source.startswith("file://")
            ):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)
        elif isinstance(source, Path):
            return self.convert_local(source, **kwargs)

    def convert_local(
        self, path: Union[str, Path], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        if isinstance(path, Path):
            path = str(path)
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)

        for g in self._guess_ext_magic(path):
            self._append_ext(extensions, g)

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(
        self, stream: Any, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(
        self, url: str, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        response = self._requests_session.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, url=response.url, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(
        self, local_path: str, extensions: List[Union[str, None]], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""

        # Create a copy of the page_converters list, sorted by priority.
        # We do this with each call to _convert because the priority of converters may change between calls.
        # The sort is guaranteed to be stable, so converters with the same priority will remain in the same order.
        sorted_converters = sorted(self._page_converters, key=lambda x: x.priority)

        for ext in extensions + [None]:  # Try last with no extension
            for converter in sorted_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                # Copy any additional global options
                if "llm_client" not in _kwargs and self._llm_client is not None:
                    _kwargs["llm_client"] = self._llm_client

                if "llm_model" not in _kwargs and self._llm_model is not None:
                    _kwargs["llm_model"] = self._llm_model

                if "style_map" not in _kwargs and self._style_map is not None:
                    _kwargs["style_map"] = self._style_map

                if "exiftool_path" not in _kwargs and self._exiftool_path is not None:
                    _kwargs["exiftool_path"] = self._exiftool_path

                # Add the list of converters for nested processing
                _kwargs["_parent_converters"] = self._page_converters

                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join(
                        [line.rstrip() for line in re.split(r"\r?\n", res.text_content)]
                    )
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)

            # Fix for: https://github.com/microsoft/markitdown/issues/222
            # If there are no guesses, then try again after trimming leading ASCII whitespaces.
            # ASCII whitespace characters are those byte values in the sequence b' \t\n\r\x0b\f'
            # (space, tab, newline, carriage return, vertical tab, form feed).
            if len(guesses) == 0:
                with open(path, "rb") as file:
                    while True:
                        char = file.read(1)
                        if not char:  # End of file
                            break
                        if not char.isspace():
                            file.seek(file.tell() - 1)
                            break
                    try:
                        guesses = puremagic.magic_stream(file)
                    except puremagic.main.PureError:
                        pass

            extensions = list()
            for g in guesses:
                ext = g.extension.strip()
                if len(ext) > 0:
                    if not ext.startswith("."):
                        ext = "." + ext
                    if ext not in extensions:
                        extensions.append(ext)
            return extensions
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return []

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
